# -*- coding: utf-8 -*-
"""Deck GNLink -> Edge: sinergias que a combinacao destrava. Externo, voz GNLink."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

NAVY=RGBColor(0x15,0x33,0x5B); TEAL=RGBColor(0x34,0xBB,0xAC); ORANGE=RGBColor(0xF2,0x7C,0x25)
DARK=RGBColor(0x22,0x2B,0x38); GRAY=RGBColor(0x6B,0x7A,0x8C); LIGHT=RGBColor(0xF1,0xF5,0xF7)
LINE=RGBColor(0xCF,0xDA,0xE1); WHITE=RGBColor(0xFF,0xFF,0xFF)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def kb(p):
    pPr=p._p.get_or_add_pPr(); pPr.set("marL","0"); pPr.set("indent","0")
    for e in pPr.findall(qn("a:buNone"))+pPr.findall(qn("a:buChar"))+pPr.findall(qn("a:buAutoNum")): pPr.remove(e)
    pPr.append(etree.SubElement(pPr,qn("a:buNone")))
def T(sl,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=2):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=tf.margin_right=Inches(0.03); tf.margin_top=tf.margin_bottom=0; tf.vertical_anchor=anchor
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align; kb(p); p.space_after=Pt(sp); p.space_before=Pt(0)
        for (s,sz,c,b,*it) in line:
            r=p.add_run(); r.text=s; r.font.name="Calibri"; r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=b
            if it and it[0]: r.font.italic=True
    return tb
def R(sl,l,t,w,h,color,ln=None,dash=False):
    sp=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if color is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=color
    if ln is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=ln; sp.line.width=Pt(1)
        if dash:
            lnEl=sp.line._get_or_add_ln(); d=lnEl.find(qn("a:prstDash"))
            if d is None: d=etree.SubElement(lnEl,qn("a:prstDash"))
            d.set("val","dash")
    sp.shadow.inherit=False; sp.text_frame.paragraphs[0].text=""; return sp
def wordmark(sl):
    T(sl,11.55,0.42,1.7,0.3,[[("GN",13,NAVY,True),("Link",13,TEAL,True)]],align=PP_ALIGN.RIGHT)
def foot(sl,n):
    T(sl,0.6,7.06,10.5,0.25,[[("Fonte: GNLink; short list de sinergias GNLink e Edge, 2026.",8.5,GRAY,False)]])
    T(sl,12.5,7.06,0.6,0.25,[[(str(n),9,GRAY,False)]],align=PP_ALIGN.RIGHT)

KIND={"MAP":"MAPA","IMG":"IMAGEM","DIAG":"DIAGRAMA","CHART":"GRÁFICO"}
def right_block(sl,spec,l,t,w,h):
    k=spec[0]
    if k in KIND:
        R(sl,l,t,w,h,RGBColor(0xF7,0xFA,0xFB),ln=GRAY,dash=True)
        T(sl,l,t+h/2-0.62,w,0.4,[[("[ "+KIND[k]+" ]",13,NAVY,True)]],align=PP_ALIGN.CENTER)
        T(sl,l+0.35,t+h/2-0.18,w-0.7,0.9,[[(spec[1],10.5,GRAY,False)]],align=PP_ALIGN.CENTER)
        T(sl,l,t+h-0.4,w,0.3,[[("a inserir",9,GRAY,False,True)]],align=PP_ALIGN.CENTER)
    elif k=="STAT":
        R(sl,l,t,w,h,LIGHT); R(sl,l,t,0.09,h,TEAL)
        T(sl,l+0.3,t+h/2-0.7,w-0.5,0.7,[[(spec[1],28,NAVY,True)]])
        T(sl,l+0.32,t+h/2+0.0,w-0.6,1.4,[[(spec[2],12,DARK,False)]])
    elif k=="CARD":
        R(sl,l,t,w,h,LIGHT); R(sl,l,t,0.09,h,TEAL)
        T(sl,l+0.3,t+0.35,w-0.55,0.8,[[(spec[1],15,NAVY,True)]])
        T(sl,l+0.32,t+1.15,w-0.6,h-1.4,[[(spec[2],12,DARK,False)]])

def content(n,bloco,cat,title,sub,points,right):
    sl=prs.slides.add_slide(BLANK); wordmark(sl)
    T(sl,0.6,0.5,10,0.3,[[("BLOCO "+bloco+"  ·  "+cat.upper(),10.5,TEAL,True)]])
    T(sl,0.6,0.86,12.0,0.95,[[(title,21,NAVY,True)]])
    T(sl,0.6,1.86,12.0,0.35,[[(sub,12,GRAY,False)]])
    pr=[[("•  ",12,TEAL,True),(pt,12,DARK,False)] for pt in points]
    T(sl,0.6,2.5,6.75,4.2,pr,sp=9)
    right_block(sl,right,7.7,2.5,5.05,4.05)
    foot(sl,n)

SYN=[
("A","Comercial","A combinação alcança 40,7% do território nacional, com praças fora do raio de Santos",
 "Cobertura territorial da plataforma combinada",
 ["As 3 plantas da GNLink já alcançam 31,1% do território brasileiro; com os projetos assinados, 40,7%.",
  "Nordeste, Sul, Paraguai e Argentina complementam o Sudeste atendido pelo terminal de Santos.",
  "É capacidade instalada em regiões onde o alcance rodoviário a partir de Santos é antieconômico."],
 ("MAP","Alcance territorial: plantas GNLink (NE, Sul, Argentina) e terminal Edge (Santos)")),
("A","Comercial","Três linhas de receita ampliam a base de clientes atendível, hoje concentrada em GNL",
 "Portfólio de produtos da GNLink",
 ["GNL para o cliente de médio e grande porte, a mesma vertente do terminal de Santos.",
  "GNC, com cadeia cerca de 3x mais eficiente em capital, atende bem o cliente pequeno.",
  "Remuneração sobre infraestrutura: receita de serviço, sem exposição direta à commodity."],
 ("CARD","3 linhas de receita","GNL, GNC e remuneração sobre infraestrutura")),
("A","Comercial","Contratos ativos cobrem clientes de pequeno, médio e grande consumo",
 "Amplitude da carteira por porte de cliente",
 ["A carteira vai do pequeno consumidor, atendido por GNC, ao grande industrial em GNL.",
  "Flexibilidade para dimensionar a solução ao perfil de cada demanda.",
  "Menor dependência de um único porte de consumo."],
 ("CARD","Do pequeno ao grande","Toda a curva de consumo coberta")),
("A","Comercial","Track record com indústria, postos e chamadas públicas de redes locais",
 "Tipos de cliente já contratados",
 ["Clientes industriais de diversos setores da economia.",
  "Postos de combustível, capturando a demanda de transporte.",
  "Vitórias em chamadas públicas para fornecimento a redes locais de distribuidoras."],
 ("CARD","3 tipos de cliente","Indústria, postos e redes locais")),
("A","Comercial","Presença em múltiplos setores, com leitura de consumo e desafios de cada um",
 "Segmentos industriais atendidos",
 ["Projetos em vários setores da economia.",
  "Conhecimento das particularidades de consumo e dos desafios de cada demanda.",
  "Base para replicar soluções entre segmentos."],
 ("CARD","Multissetorial","Leitura de consumo e desafios por segmento")),
("A","Comercial","Time comercial com raízes em GLP e óleo combustível, os substitutos a vencer",
 "Origem e foco do time comercial",
 ["Experiência com GLP e óleo combustível, os principais substitutos no cliente industrial.",
  "GLP e óleo combustível representam 81% do pipeline da GNLink.",
  "Conhecimento prático de como deslocar o combustível incumbente no cliente."],
 ("STAT","81%","do pipeline em GLP e óleo combustível, os substitutos a vencer")),
("A","Comercial","A GNLink já entrega 0,28 MMm³/d faturando, reforçando a execução que o mercado acompanha",
 "Volume já entregue vs. curva de rampa do off-grid",
 ["O SS-LNG off-grid tende a 50% do lucro bruto da Edge até 2030, saindo de 0,15 para 3,80 MMm³/d em 3 anos.",
  "A GNLink já fatura 0,28 MMm³/d e tem 0,4 MMm³/d assinado.",
  "Conversão de cliente comprovada reforça o fator de execução que o mercado mais acompanha."],
 ("CHART","Curva de rampa do off-grid vs. volume já entregue pela GNLink")),
("A","Engenharia","Engenharia própria projeta equipamentos nacionais mais baratos, cortando orçamento e prazo",
 "Capacidade de engenharia da GNLink",
 ["O time desenhou instalações e equipamentos bem mais baratos que as soluções de mercado.",
  "Soluções nacionais reduzem o custo e o prazo de implantação no cliente.",
  "É a base da fábrica de regás in-house (ver bloco B)."],
 ("CARD","Nacional e mais barato","Menos orçamento, menos prazo de implantação")),
("A","Comercial","Contratos com postos já capturam a demanda de gás no transporte, mercado de grande potencial",
 "Presença no segmento de transporte",
 ["Vários contratos assinados com postos de combustível.",
  "Captura inicial da demanda por gás no transporte rodoviário e nos corredores azuis.",
  "Segmento de mercado potencial enorme, ainda pouco penetrado."],
 ("CARD","Transporte","Postos e corredores azuis, mercado de grande potencial")),
("A","Operacional","Três plantas operacionais adicionam redundância ao terminal único de Santos",
 "Resiliência de fornecimento da plataforma combinada",
 ["Um ativo único carrega risco de não fornecimento durante uma parada.",
  "As 3 plantas da GNLink e o terminal da Edge podem se cobrir mutuamente (backup cruzado).",
  "Maior confiabilidade percebida pelo cliente off-grid."],
 ("DIAG","Backup cruzado: terminal Edge e as 3 plantas GNLink se cobrindo")),
("B","Operacional","Capacidade adicional pronta em RN e Itabuna/BA, com licença, terreno e molécula assegurados",
 "Expansões de baixo risco já mapeadas",
 ["RN: licença e terreno já disponíveis para ampliar volume.",
  "Itabuna/BA: terreno ao lado da planta atual.",
  "Sem risco de falta de molécula nas duas praças."],
 ("CARD","RN e Itabuna/BA","Licença, terreno e molécula já assegurados")),
("B","Estratégico","Capacidade instalada e diversificada complementa o plano de expansão da Edge",
 "Ativo pronto vs. capacidade a construir",
 ["O plano da Edge amplia de forma expressiva a capacidade dentro do raio de Santos.",
  "A GNLink já tem ativo instalado, operando e com mercado diversificado.",
  "Somar capacidade pronta e diversificada reduz o risco de execução da curva de crescimento."],
 ("CARD","Ativo pronto","Complementa a capacidade que o plano ainda vai construir")),
("B","Molécula","Molécula do TRSP alimenta Itabuna, e as fontes onshore da GNLink reforçam o mix da Edge",
 "Fluxos de molécula da entidade combinada, bidirecionais",
 ["Molécula de mercado livre a partir do TRSP pode alimentar a planta de Itabuna/BA, conectada ao grid.",
  "As fontes onshore da GNLink fortalecem o mix de suprimento da Edge.",
  "Mais fontes e mais flexibilidade de despacho, por cliente e por mês."],
 ("DIAG","Fluxos de molécula: TRSP, plantas GNLink e campos onshore")),
("B","Molécula","Opção de assumir a molécula própria no Sul, via o campo da Tradener",
 "Alavanca adicional de suprimento, a avaliar",
 ["A compra do campo da Tradener tornaria a GNLink dona da molécula no Sul.",
  "Fortalece o mix e reduz a dependência de terceiros na região.",
  "Ponto a aprofundar em conjunto na avaliação."],
 ("CARD","A avaliar em conjunto","Campo da Tradener: molécula própria no Sul")),
("B","Logística","Carretas ociosas da Edge podem servir as plantas da GNLink e reduzir o OpEx logístico",
 "Otimização da frota combinada",
 ["A Edge já adquiriu carretas para a sua operação.",
  "A capacidade ociosa pode atender rotas das plantas da GNLink.",
  "Redução do custo logístico atual, sem novo investimento."],
 ("CARD","Frota combinada","Carretas ociosas reduzem o OpEx logístico")),
("B","Financeiro","A robustez financeira da Compass otimiza a estrutura de capital da GNLink",
 "Custo de capital da GNLink vs. porte da Compass",
 ["A GNLink capta hoje entre 10,7% e 14,25% sobre cerca de R$400 mi de dívida.",
  "A Compass é listada, com ND/EBITDA de 2,1x e acesso amplo ao mercado.",
  "Cada 100 bps de redução equivale a cerca de R$4 mi por ano."],
 ("STAT","~R$4 mi/ano","por 100 bps de redução no custo da dívida (~R$400 mi)")),
("B","CapEx","A única fábrica de regás nacional corta capex de importação e acelera a implantação",
 "Regás in-house vs. importada",
 ["Uma regás de 200 m³/h custa cerca de R$1 mi pela GNLink, contra cerca de R$3 mi importada (3x).",
  "Viabiliza regás de pequeno porte, ajustada ao perfil de cliente e à capilaridade das plantas.",
  "Ordem de R$160 mi de capex de importação evitado, e implantação mais rápida no cliente."],
 ("IMG","Unidade de regás in-house da GNLink")),
("B","Cross-segmento","A GNLink expande a malha das 8 distribuidoras da Compass, o núcleo de ~65% do equity",
 "Potencial de redes locais em concessões Compass",
 ["A GNLink leva o gás até a cidade e a distribuidora constrói só a malha interna, sem puxar ramal da capital.",
  "Todo o potencial de rede local do Sul (433 mil m³/d) está em concessões da própria Compass.",
  "Com o Nordeste, o potencial mapeado chega a 916 mil m³/d."],
 ("MAP","Redes locais potenciais nas concessões das distribuidoras Compass")),
("B","Regulatório","O músculo regulatório da Compass acelera o licenciamento das plantas da GNLink",
 "Capacidade institucional aplicada aos gargalos atuais",
 ["Gargalos datados de licenciamento (ANP, Coelba, PetroRecôncavo) em resolução.",
  "A estrutura regulatória da Compass encurta prazos e reduz incerteza.",
  "Destrava volume nas plantas existentes mais rápido."],
 ("CARD","Licenciamento","Prazos menores nos gargalos hoje datados")),
("B","Capital","Equity listado da Compass como moeda de aquisição, com escudo fiscal de prejuízos acumulados",
 "Flexibilidade de estrutura da transação",
 ["O equity listado da Compass pode compor a estrutura da combinação.",
  "Prejuízos fiscais acumulados formam um escudo aproveitável.",
  "Amplia as opções de desenho da operação."],
 ("CARD","Flexibilidade","Equity listado como moeda e escudo fiscal")),
("B","Novos projetos","A Argentina adiciona 300k m³/d de EBITDA sem capital empregado, com capex do supridor",
 "Crescimento a capital empregado zero",
 ["Na Argentina o capex é do supridor da molécula: a GNLink não emprega capital.",
  "A margem (cerca de R$1,13/m³) supera a das praças brasileiras, pela molécula mais barata.",
  "Cerca de 300k m³/d e R$100 mi de EBITDA, com início previsto para o 1T28."],
 ("MAP","Paso de los Libres e hub Uruguaiana: alcance no Cone Sul")),
("B","Novos projetos","O projeto Eneva adiciona crescimento contratado a capex leve, já em 2027",
 "Pipeline contratado de curto prazo",
 ["Take assinado com a Eneva.",
  "Capex leve e entrada em 2027.",
  "Crescimento incremental já contratado, somando à curva."],
 ("CARD","Entrada em 2027","Take assinado, capex leve")),
]

SHORT_A=["Extensão geográfica","Diversidade de produto","Diversidade de tamanho de cliente","Diversidade de tipos de cliente",
 "Diversidade de segmentos industriais","Expertise comercial","Execução comprovada","Expertise de engenharia",
 "Expertise com transporte","Backup e resiliência"]
SHORT_B=["Crescimento das plantas atuais","Redução do risco de capex","Molécula, fontes e flexibilidade",
 "Molécula própria no Sul (a avaliar)","Otimização logística","Estrutura de capital","Fábrica de regás in-house",
 "Redes locais das distribuidoras","Aceleração regulatória","Equity e escudo fiscal","Argentina, capex do supridor","Projeto Eneva"]

# CAPA
c=prs.slides.add_slide(BLANK)
R(c,0,0,13.333,7.5,NAVY); R(c,0,6.25,13.333,0.12,TEAL)
T(c,0.9,2.2,11,0.6,[[("GN",16,WHITE,True),("Link",16,TEAL,True),("   ×   Edge",16,WHITE,True)]])
T(c,0.9,2.75,11.5,1.2,[[("As sinergias que a combinação destrava",34,WHITE,True)]])
T(c,0.9,3.95,11,0.5,[[("Uma leitura das complementaridades entre as duas plataformas de gás off-grid",14,TEAL,False)]])
T(c,0.9,6.5,11,0.35,[[("Documento para discussão  ·  2026",11,RGBColor(0xB8,0xC6,0xD2),False)]])

# RESUMO
s=prs.slides.add_slide(BLANK); wordmark(s)
T(s,0.6,0.5,10,0.3,[[("VISÃO GERAL",10.5,TEAL,True)]])
T(s,0.6,0.86,12.2,0.6,[[("Vinte e duas sinergias, em duas frentes complementares",21,NAVY,True)]])
T(s,0.6,1.8,12.2,0.35,[[("O que a GNLink traz para a plataforma, e o que a combinação das duas destrava",12,GRAY,False)]])
def col(x,head,items,pref):
    T(s,x,2.4,5.9,0.35,[[(head,13,NAVY,True)]]); R(s,x,2.78,5.9,0.03,TEAL); y=2.98
    for i,it in enumerate(items):
        T(s,x,y,5.9,0.36,[[(pref+str(i+1)+"  ",11,TEAL,True),(it,11,DARK,False)]]); y+=0.365
col(0.6,"A  ·  O que a GNLink traz",SHORT_A,"A")
col(6.9,"B  ·  O que a combinação destrava",SHORT_B,"B")
foot(s,2)

# 22 SLIDES
n=3
for (bloco,cat,title,sub,points,right) in SYN:
    content(n,bloco,cat,title,sub,points,right); n+=1

# FECHAMENTO
f=prs.slides.add_slide(BLANK); wordmark(f)
T(f,0.6,0.5,10,0.3,[[("SÍNTESE",10.5,TEAL,True)]])
T(f,0.6,0.86,12.2,0.9,[[("Uma plataforma combinada mais ampla, mais resiliente e mais rápida",21,NAVY,True)]])
cards=[("Mais ampla","Geografia, produtos e tipos de cliente que multiplicam o mercado endereçável da plataforma."),
       ("Mais resiliente","Backup cruzado, mix de molécula e a robustez financeira da Compass reduzem os riscos do conjunto."),
       ("Mais rápida","Ativo já instalado, regás in-house e músculo regulatório aceleram a curva de crescimento.")]
for i,(h,b) in enumerate(cards):
    x=0.6+i*4.15; R(f,x,2.2,3.9,2.1,LIGHT); R(f,x,2.2,3.9,0.1,TEAL)
    T(f,x+0.28,2.5,3.4,0.5,[[(h,16,NAVY,True)]]); T(f,x+0.3,3.05,3.35,1.2,[[(b,12,DARK,False)]])
T(f,0.6,4.75,12,0.35,[[("PRÓXIMOS PASSOS",10.5,TEAL,True)]])
R(f,0.6,5.1,12.15,1.35,RGBColor(0xF7,0xFA,0xFB),ln=GRAY,dash=True)
T(f,0.6,5.62,12.15,0.4,[[("[ a definir em conjunto: agenda de diligência e quantificação das sinergias prioritárias ]",12,GRAY,False)]],align=PP_ALIGN.CENTER)
foot(f,n)

prs.save("GNLink-Sinergias-Edge.pptx")
print("OK. Total de slides:",len(prs.slides._sldIdLst))
