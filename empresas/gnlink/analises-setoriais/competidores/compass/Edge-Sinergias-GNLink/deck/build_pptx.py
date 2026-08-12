# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

OUT = sys.argv[1]

NAVY=RGBColor(0x00,0x2D,0x5C); MED=RGBColor(0x28,0x47,0x7E); ORANGE=RGBColor(0xEF,0x76,0x22)
DKOR=RGBColor(0xC5,0x5A,0x17); SAGE=RGBColor(0x77,0xA2,0x8A); MINT=RGBColor(0xA8,0xD2,0xBB)
GRAY=RGBColor(0x7F,0x7F,0x7F); TEALD=RGBColor(0x5A,0x9B,0x87); WHITE=RGBColor(0xFF,0xFF,0xFF)
OFF=RGBColor(0xF2,0xF4,0xF7); DK=RGBColor(0x1E,0x1E,0x1E); LINE=RGBColor(0xD9,0xE0,0xE8)
HLROW=RGBColor(0xEA,0xF1,0xEC); HLBOX=RGBColor(0xEA,0xF1,0xEC); PEACH=RGBColor(0xFB,0xE7,0xD6); RED=RGBColor(0xB9,0x43,0x3A)

prs=Presentation(OUT)
COV=[l for l in prs.slide_layouts if l.name=="Template Branco"][0]
CON=[l for l in prs.slide_layouts if l.name=="Escolhido: Sem Subtexto; Sem Fonte"][0]
CLO=[l for l in prs.slide_layouts if l.name=="7_Título e conteúdo"][0]

xml_slides=prs.slides._sldIdLst
for sid in list(xml_slides)[1:]:
    prs.part.drop_rel(sid.get(qn("r:id"))); xml_slides.remove(sid)

def kill_bul(p):
    pPr=p._p.get_or_add_pPr(); pPr.set("marL","0"); pPr.set("indent","0")
    for e in pPr.findall(qn("a:buNone"))+pPr.findall(qn("a:buChar"))+pPr.findall(qn("a:buAutoNum")): pPr.remove(e)
    pPr.append(etree.SubElement(pPr,qn("a:buNone")))
def strip_shape_bul(shape):
    if not shape.has_text_frame: return
    ls=shape.text_frame._txBody.find(qn("a:lstStyle"))
    if ls is None: return
    for lp in list(ls):
        for bu in lp.findall(qn("a:buChar"))+lp.findall(qn("a:buAutoNum"))+lp.findall(qn("a:buFont")): lp.remove(bu)
        lp.append(etree.SubElement(lp,qn("a:buNone"))); lp.set("marL","0"); lp.set("indent","0")
def fill_topbar(slide, section, title, section_size=None, title_size=15):
    for shp in slide.shapes:
        if not (shp.has_text_frame and "Placeholder" in shp.name): continue
        try: idx=shp.placeholder_format.idx
        except: continue
        if idx==13:
            strip_shape_bul(shp); tf=shp.text_frame; tf.clear(); p=tf.paragraphs[0]; kill_bul(p)
            r=p.add_run(); r.text=section; r.font.color.rgb=NAVY; r.font.name="Calibri"; r.font.bold=True
            if section_size: r.font.size=Pt(section_size)
        elif idx==14:
            strip_shape_bul(shp); tf=shp.text_frame; tf.clear(); p=tf.paragraphs[0]; kill_bul(p)
            r=p.add_run(); r.text=title; r.font.color.rgb=NAVY; r.font.name="Calibri"; r.font.size=Pt(title_size)
def num(slide,n):
    tb=slide.shapes.add_textbox(Inches(12.5),Inches(7.12),Inches(0.6),Inches(0.25)); tf=tb.text_frame
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    r=p.add_run(); r.text=str(n); r.font.name="Calibri"; r.font.size=Pt(9); r.font.color.rgb=GRAY
def footer(slide, txt_):
    tb=slide.shapes.add_textbox(Inches(0.55),Inches(7.10),Inches(11.6),Inches(0.25)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0; p=tf.paragraphs[0]
    r=p.add_run(); r.text="Fonte: "+txt_; r.font.name="Calibri"; r.font.size=Pt(8); r.font.color.rgb=GRAY
def setfill(shape,color):
    shape.fill.solid(); shape.fill.fore_color.rgb=color; shape.line.fill.background()
def txt(slide,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,space=None):
    tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.margin_left=tf.margin_right=Inches(0.02); tf.margin_top=tf.margin_bottom=0; tf.vertical_anchor=anchor
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        if space: p.space_after=Pt(space)
        kill_bul(p)
        for (s,sz,c,b) in line:
            r=p.add_run(); r.text=s; r.font.name="Calibri"; r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=b
    return tb
def rect(slide,l,t,w,h,color,shape=MSO_SHAPE.RECTANGLE):
    sp=slide.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); setfill(sp,color); sp.shadow.inherit=False
    sp.text_frame.paragraphs[0].text=""; return sp
def stat(slide,l,t,w,number,label,ncolor=NAVY,nsize=27):
    txt(slide,l,t,w,0.5,[[(number,nsize,ncolor,True)]]); txt(slide,l,t+0.48,w,0.62,[[(label,10,GRAY,False)]])
def content_slide(section,title,section_size=None,title_size=15,pagenum=None):
    s=prs.slides.add_slide(CON)
    for shp in list(s.shapes):
        if shp.has_text_frame and "Placeholder" in shp.name:
            try: idx=shp.placeholder_format.idx
            except: idx=None
            if idx not in (13,14): shp._element.getparent().remove(shp._element)
    fill_topbar(s,section,title,section_size,title_size)
    if pagenum is not None: num(s,pagenum)
    return s
def divider(numlbl,title,pagenum):
    s=prs.slides.add_slide(CON)
    for shp in list(s.shapes):
        if shp.has_text_frame and "Placeholder" in shp.name: shp._element.getparent().remove(shp._element)
    rect(s,-0.1,-0.1,13.6,7.7,NAVY); rect(s,0.9,3.05,1.4,0.09,ORANGE)
    txt(s,0.9,2.1,11,0.9,[[(numlbl,54,ORANGE,True)]]); txt(s,0.9,3.25,11.5,1.4,[[(title,30,WHITE,True)]])
    num(s,pagenum); return s

NOSTYLE="{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
def set_border(cell,edges,color=LINE,w=6350):
    tcPr=cell._tc.get_or_add_tcPr()
    for ed in edges:
        tag=qn("a:ln"+ed)
        for ex in tcPr.findall(tag): tcPr.remove(ex)
        ln=etree.SubElement(tcPr,tag); ln.set("w",str(w)); ln.set("cap","flat")
        fill=etree.SubElement(ln,qn("a:solidFill")); clr=etree.SubElement(fill,qn("a:srgbClr")); clr.set("val","%02X%02X%02X"%(color[0],color[1],color[2]))
def table(slide,data,l,t,w,colw,fs=9.5,header=True,hl_last_col=False,hl_rows=None,rowh=0.30,hcolor=NAVY):
    rows=len(data); cols=len(data[0])
    gtbl=slide.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(w),Inches(rowh*rows)); tbl=gtbl.table
    tbl.first_row=False; tbl.horz_banding=False; tbl.last_row=False; tbl.first_col=False; tbl.last_col=False
    tbl._tbl.tblPr.set("firstRow","0"); tbl._tbl.tblPr.set("bandRow","0")
    st=tbl._tbl.tblPr.find(qn("a:tableStyleId"))
    if st is None: st=etree.SubElement(tbl._tbl.tblPr,qn("a:tableStyleId"))
    st.text=NOSTYLE
    tot=sum(colw)
    for j,cw in enumerate(colw): tbl.columns[j].width=Emu(int(Inches(w)*cw/tot))
    hl_rows=hl_rows or []
    for i,row in enumerate(data):
        tbl.rows[i].height=Inches(rowh)
        for j,val in enumerate(row):
            cell=tbl.cell(i,j)
            cell.margin_left=Inches(0.07); cell.margin_right=Inches(0.05); cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            ishead=header and i==0
            if ishead: cell.fill.solid(); cell.fill.fore_color.rgb=hcolor
            elif i in hl_rows: cell.fill.solid(); cell.fill.fore_color.rgb=HLROW
            else: cell.fill.solid(); cell.fill.fore_color.rgb=WHITE if i%2==1 else OFF
            set_border(cell,["B"],LINE,6350)
            tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; kill_bul(p)
            p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
            bold = ishead or (hl_last_col and j==cols-1) or (i in hl_rows)
            r=p.add_run(); r.text=str(val); r.font.name="Calibri"; r.font.size=Pt(fs if not ishead else fs-0.5); r.font.bold=bold
            r.font.color.rgb=WHITE if ishead else (NAVY if bold else DK)
    return gtbl

# ================= SLIDES =================
cov=prs.slides[0]
for shp in cov.shapes:
    if not shp.has_text_frame: continue
    t=shp.text_frame.text
    if "Título" in t:
        strip_shape_bul(shp); tf=shp.text_frame; tf.clear(); p=tf.paragraphs[0]; kill_bul(p)
        r=p.add_run(); r.text="GNLink × Edge"; r.font.name="Calibri"; r.font.bold=True; r.font.size=Pt(32); r.font.color.rgb=WHITE
    elif "Dia" in t or "Mês" in t or "Ano" in t:
        strip_shape_bul(shp); tf=shp.text_frame; tf.clear(); p=tf.paragraphs[0]; kill_bul(p)
        r=p.add_run(); r.text="Tese de combinação estratégica"; r.font.name="Calibri"; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=MINT
        p2=tf.add_paragraph(); kill_bul(p2)
        r2=p2.add_run(); r2.text="Julho 2026  ·  Confidencial"; r2.font.name="Calibri"; r2.font.size=Pt(11); r2.font.color.rgb=WHITE

pg=1
# ---- 2 EXEC SUMMARY ----
pg+=1; s=content_slide("Sumário executivo","A Edge sondou a GNLink; a combinação é complementar e a régua situa o ativo em R$1,4 a 1,6 bilhão",title_size=13.5,pagenum=pg)
for i,(lab,body,col) in enumerate([
 ("Situação","A Edge (Compass) sondou a GNLink para M&A. As duas competem em GNL off-grid, mas com modelos e geografias complementares.",MED),
 ("Complicação","A Fase 2 da Edge exige R$1,5 bi para crescer 25x na vertical que vale metade do seu lucro bruto; o mercado aponta a execução como risco número um.",ORANGE),
 ("Resolução","A GNLink entrega hoje 686 mil m³/d e R$230 mi de EBITDA, na margem que o próprio modelo da Edge projeta, com capex incremental zero.",SAGE)]):
    y=1.25+i*0.62; rect(s,0.55,y,0.12,0.5,col)
    txt(s,0.78,y,1.55,0.5,[[(lab,12,col,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,2.4,y,10.35,0.55,[[(body,11,DK,False)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,0.55,3.35,12.2,0.02,LINE)
mets=[("3","plantas operando (PR, BA, RN)"),("280","mil m³/d, 79% já contratado"),("686","mil m³/d com o pipeline"),("R$230 mi","EBITDA na curva madura"),("R$1,37 bi","referência de valor (EBITDA 2028, 10x)")]
for i,(n_,l_) in enumerate(mets):
    stat(s,0.55+i*2.44,3.6,2.35,n_,l_)
cards=[("Execução comprovada","3 plantas e 15 contratos faturando na vertical que o mercado mais valoriza."),
       ("Crescimento sem balanço","Mais 406 mil m³/d contratados a capex zero (Eneva e Argentina)."),
       ("Sinergias reais","G&A de R$210 a 360 mi, refinanciamento e redes locais das 7 distribuidoras."),
       ("Valor destravado","A GNLink desloca cerca de R$1 bi no valuation da Edge, perto de +3,5% no preço-alvo.")]
for i,(h,b) in enumerate(cards):
    x=0.55+i*3.06; rect(s,x,4.85,2.9,1.75,OFF); rect(s,x,4.85,2.9,0.08,ORANGE if i==3 else SAGE)
    txt(s,x+0.15,5.02,2.6,0.5,[[(h,12,NAVY,True)]]); txt(s,x+0.15,5.45,2.62,1.1,[[(b,10,DK,False)]])
footer(s,"Análise Lorinvest; equity research (BofA, Citi, Itaú BBA, BTG, Bradesco BBI); dados da companhia")

pg+=1; divider("01","Contexto e tese",pg)

# ---- 4 SNAPSHOT ----
pg+=1; s=content_slide("Contexto","A GNLink lidera o gás onde o gasoduto não chega: liquefaz no país e leva o GNL por caminhão e barcaça",title_size=13.5,pagenum=pg)
txt(s,0.55,1.25,6.2,0.35,[[("Modelo de negócio",13,NAVY,True)]])
flow=[("Molécula nacional","gás de gasoduto"),("Liquefação","plantas próprias"),("Logística","carreta, GNC, barcaça"),("Cliente off-grid","indústria, frota")]
for i,(a,b) in enumerate(flow):
    x=0.55+i*1.6; rect(s,x,1.7,1.45,0.85,OFF)
    txt(s,x+0.05,1.78,1.35,0.8,[[(a,10,NAVY,True)],[(b,8.5,GRAY,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if i<3: txt(s,x+1.45,1.7,0.16,0.85,[[("›",17,SAGE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.55,2.85,6.2,0.35,[[("Posição",13,NAVY,True)]])
for i,li in enumerate(["Base nacional, enquanto os concorrentes atuam por região","Engenharia de regaseificação in-house, única no país","Time comercial ex-GLP, especializado em conversão de cliente","Duas modalidades: GNL para o cliente grande e GNC para o pequeno"]):
    y=3.25+i*0.5; rect(s,0.6,y+0.05,0.1,0.1,SAGE); txt(s,0.85,y,6.0,0.45,[[(li,11,DK,False)]])
txt(s,7.1,1.25,5.6,0.35,[[("Indicadores",13,NAVY,True)]])
kpis=[("51","colaboradores"),("R$1 bi+","em contratos de venda"),("R$0,78/m³","margem EBITDA das plantas"),("2022","constituição, 3 de 6 projetos")]
for i,(n_,l_) in enumerate(kpis):
    x=7.1+(i%2)*2.9; y=1.75+(i//2)*1.35; rect(s,x,y,2.7,1.1,OFF)
    txt(s,x+0.18,y+0.14,2.4,0.55,[[(n_,26,NAVY,True)]]); txt(s,x+0.18,y+0.68,2.45,0.4,[[(l_,10,GRAY,False)]])
rect(s,7.1,4.6,5.6,1.9,HLBOX); rect(s,7.1,4.6,5.6,0.08,ORANGE)
txt(s,7.28,4.78,5.3,0.4,[[("Alcance completo, em duas tecnologias",12,NAVY,True)]])
txt(s,7.28,5.2,5.3,1.2,[[("GNL para o cliente médio-grande (média de 30,5 mil m³/d) e GNC, com capex 3x menor, para o cliente pequeno, faixa que o modelo de terminal não atende.",11,DK,False)]])
footer(s,"Histórico GNLink 2026; Investment Presentation 2025; dados da companhia")

# ---- 5 WHY EDGE ----
pg+=1; s=content_slide("Tese","A GNLink já opera a vertical que responde por metade do lucro bruto da Edge e que precisa crescer 25x em três anos",title_size=12.5,pagenum=pg)
txt(s,0.55,1.25,6.1,0.35,[[("Small Scale LNG dentro da Edge",13,NAVY,True)]])
table(s,[["Ano","2027","2028","2029","2030"],
         ["Volume (mi m³/d)","0,15","1,37","2,58","3,80"],
         ["Margem (R$/m³)","0,80","","","0,86"],
         ["Lucro bruto (R$ mi)","44","410","792","1.195"]],
      0.55,1.7,6.1,[2.1,1,1,1,1],fs=10,rowh=0.42)
txt(s,0.55,3.7,6.1,1.4,[[("O off-grid vira metade do lucro bruto da Edge até 2030, saindo de quase zero. O ritmo de conversão de clientes é o fator que os cinco bancos mais acompanham.",11,DK,False)]])
for i,(n_,l_,c) in enumerate([("25x","de crescimento exigido em três anos",ORANGE),("~0,25 mi m³/d","já entregues pela GNLink, faturando",SAGE),("1,9x","o que a Edge projeta operar em 2027",NAVY)]):
    y=1.4+i*1.7; rect(s,7.0,y,5.7,1.45,OFF); rect(s,7.0,y,0.1,1.45,c)
    txt(s,7.25,y+0.12,5.4,0.7,[[(n_,30,c,True)]]); txt(s,7.25,y+0.85,5.4,0.5,[[(l_,11.5,DK,False)]])
footer(s,"Equity research (Itaú BBA, BTG, Citi, BofA, Bradesco BBI) 2026")

pg+=1; divider("02","A oferta",pg)

# ---- 7 CAPEX ----
pg+=1; s=content_slide("A oferta","GNLink e Edge têm eficiência de capital comprovada equivalente; o contraste real é com a Fase 2, não comprovada",title_size=12.5,pagenum=pg)
txt(s,0.55,1.25,12,0.3,[[("Eficiência de capital, R$ por m³/d de capacidade (ex-regás)",11,GRAY,False)]])
table(s,[["","CapEx","Capacidade","R$/m³/d","Natureza"],
         ["Edge, Fase 1","R$450 mi","400 mil m³/d","R$1.125","entregue e operando"],
         ["Edge, Fase 2","R$1.500 mi","+3.400 mil m³/d","R$441","previsão não comprovada"],
         ["GNLink, 3 plantas","R$294,3 mi","280.445 m³/d","R$1.049","entregue e operando"]],
      0.55,1.65,12.2,[2.4,1.4,1.7,1.3,2.6],fs=11,rowh=0.5,hl_rows=[3])
txt(s,0.55,4.0,12,1.0,[[("A GNLink constrói a R$1.049/m³/d, 6,8% melhor que a Fase 1 da Edge, na mesma base. A Fase 2 pressupõe a intensidade de capex ",11,DK,False),("cair 61%",11,ORANGE,True),(", patamar que nenhuma das duas jamais demonstrou.",11,DK,False)]])
txt(s,0.55,4.9,12,0.3,[[("Curva de capacidade da Edge, mil m³/d",11,GRAY,False)]])
table(s,[["","2026 a 2027","2028","2029","2030+"],
         ["Edge","150","1.367","2.583","3.800"]],
      0.55,5.25,7.5,[1.4,1.5,1.1,1.1,1.1],fs=10.5,rowh=0.4)
txt(s,8.3,5.2,4.4,1.3,[[("A Fase 1 opera hoje 150 de 400 mil m³/d (38% de utilização). A GNLink, com 280 hoje e 686 com os projetos, é 1,9x a operação real da Edge em 2027.",10.5,DK,False)]])
footer(s,"Dados da companhia; ata RI Compass jan/2026; equity research 2026")

# ---- 8 CURVA ----
pg+=1; s=content_slide("A oferta","A curva completa entrega 686 mil m³/d e R$230 mi de EBITDA a capex incremental zero",title_size=13.5,pagenum=pg)
table(s,[["Bloco","Volume (mil m³/d)","EBITDA (R$ mi)","CapEx incremental"],
         ["3 plantas em operação","280,4","53,7","R$349,9 mi (afundado)"],
         ["+ G&A Holding (sinergia direta)","","+26,3",""],
         ["= Ativo visto pela Edge","280,4","80,0","zero"],
         ["+ Projeto Eneva (2027)","+106","+50","zero"],
         ["+ Argentina Fase 1 (1T28)","+150","+50","zero"],
         ["+ Argentina Fase 2 (2030)","+150","+50","zero"],
         ["Total","686,4","230,0","zero incremental"]],
      0.55,1.3,8.0,[3.0,1.6,1.4,2.0],fs=10.5,rowh=0.44,hl_rows=[3,7])
rect(s,8.9,1.3,3.85,1.75,OFF); rect(s,8.9,1.3,0.09,1.75,SAGE)
txt(s,9.1,1.46,3.5,0.4,[[("Sinergia na veia",12,NAVY,True)]])
txt(s,9.1,1.9,3.5,1.1,[[("O G&A de R$26 mi é absorvido na estrutura existente, elevando o EBITDA relevante para a Edge de R$54 para R$80 mi.",10.5,DK,False)]])
rect(s,8.9,3.25,3.85,1.75,HLBOX); rect(s,8.9,3.25,0.09,1.75,ORANGE)
txt(s,9.1,3.41,3.5,0.4,[[("Capex do supridor",12,NAVY,True)]])
txt(s,9.1,3.85,3.5,1.1,[[("Na Argentina, o CapEx é do supridor da molécula. A GNLink não emprega capital e ainda com margem superior à brasileira.",10.5,DK,False)]])
txt(s,0.55,5.55,12,0.9,[[("Diferente da Fase 2 da Edge, que exige R$1,5 bi de capital próprio com rampa a comprovar, os 406 mil m³/d de crescimento da GNLink já estão contratados e não consomem balanço.",11,NAVY,True)]])
footer(s,"Projeção da companhia; term sheets assinados (Eneva, Argentina)")

# ---- 9 ABERTURA ----
pg+=1; s=content_slide("A oferta","A abertura por planta de 2029 fecha em todas as linhas e sobrevive à due diligence",title_size=13.5,pagenum=pg)
table(s,[["GNLink 2029","PR","BA","RN","Plantas","Holding","Total"],
         ["Volume GNL (m³/d)","58.070","85.164","70.419","213.653","","213.653"],
         ["Volume GNC (m³/d)","17.496","","15.496","32.992","","32.992"],
         ["Serviço sem molécula","","27.800","6.000","33.800","","33.800"],
         ["Volume total (m³/d)","75.566","112.964","91.915","280.445","","280.445"],
         ["EBITDA (R$ mi)","25,5","25,6","28,9","80,0","(26,3)","53,7"],
         ["Margem (R$/m³)","0,92","0,62","0,86","0,78","","0,52"]],
      0.55,1.3,12.2,[2.3,1.3,1.3,1.3,1.5,1.3,1.5],fs=10,rowh=0.42,hl_rows=[4,5])
txt(s,0.55,4.6,12,0.35,[[("Duas capacidades que a Edge não tem",13,NAVY,True)]])
for i,(h,b) in enumerate([("GNC, 12% do volume","Compressão e descompressão, com capex 3x menor que a cadeia de GNL (R$543 vs R$1.514/m³/d). Atende bem o cliente pequeno."),
                          ("Serviço sem molécula, 12%","Receita de infraestrutura, sem exposição a commodity nem capital de giro de molécula.")]):
    x=0.55+i*6.2; rect(s,x,5.0,5.9,1.45,OFF); rect(s,x,5.0,5.9,0.08,SAGE)
    txt(s,x+0.15,5.16,5.6,0.4,[[(h,12,NAVY,True)]]); txt(s,x+0.15,5.58,5.62,0.85,[[(b,10.5,DK,False)]])
footer(s,"Modelo da companhia (GNLink 2029)")

# ---- 10 TRAJETORIA ----
pg+=1; s=content_slide("A oferta","O EBITDA sai do ramp-up para R$230 mi e vira positivo já em 2027",title_size=14,pagenum=pg)
txt(s,0.55,1.25,12,0.3,[[("Trajetória de EBITDA visto pela Edge, R$ mi",11,GRAY,False)]])
years=[("2026",-9.9,"-10"),("2027",97.3,"97"),("2028",172.1,"172"),("2029",180.0,"180"),("Maduro",230.0,"230")]
base_y=4.85; maxv=230; unit=2.7/maxv; cx=1.2
for i,(yr,v,disp) in enumerate(years):
    x=cx+i*2.25; h=abs(v)*unit; col=RED if v<0 else (NAVY if yr=="Maduro" else SAGE)
    if v>=0:
        rect(s,x,base_y-h,1.5,h,col); txt(s,x-0.1,base_y-h-0.3,1.7,0.3,[[("R$"+disp,12,NAVY,True)]],align=PP_ALIGN.CENTER)
    else:
        rect(s,x,base_y,1.5,h,col); txt(s,x-0.1,base_y+h+0.02,1.7,0.3,[[("R$"+disp,11,RED,True)]],align=PP_ALIGN.CENTER)
    txt(s,x-0.1,base_y+0.42,1.7,0.3,[[(yr,11,GRAY,False)]],align=PP_ALIGN.CENTER)
rect(s,1.0,base_y,11.3,0.02,GRAY)
for i,(h,b) in enumerate([("Vira positivo em 2027","As três plantas atingem EBITDA positivo por conta própria; a Eneva acelera a virada."),
                          ("79% já contratado","222 de 280 mil m³/d sob contrato; a maior parte do salto não depende de novos clientes."),
                          ("R$204 mi mais R$26 mi de G&A","O maduro visto pela Edge chega a R$230 mi, com o overhead absorvido.")]):
    x=0.55+i*4.08; rect(s,x,5.6,3.9,0.95,OFF); rect(s,x,5.6,3.9,0.07,ORANGE if i==2 else SAGE)
    txt(s,x+0.13,5.71,3.7,0.35,[[(h,11,NAVY,True)]]); txt(s,x+0.13,6.05,3.7,0.5,[[(b,9.5,DK,False)]])
footer(s,"Projeção da companhia")

# ---- 11 MARGEM ----
pg+=1; s=content_slide("A oferta","A GNLink entrega hoje a margem de R$0,80/m³ que o próprio modelo da Edge projeta para o off-grid",title_size=12.5,pagenum=pg)
table(s,[["","Volume","EBITDA","Margem (R$/m³)"],
         ["3 plantas plenas (2029)","280,4 mil m³/d","R$80 mi","0,78"],
         ["+ Eneva","106 mil m³/d","+R$50 mi","1,29"],
         ["+ Argentina F1 e F2","300 mil m³/d","+R$100 mi","0,91 cada"],
         ["Curva cheia","686,4 mil m³/d","R$230 mi","0,92"],
         ["Premissa off-grid da própria Edge","","","0,80 a 0,86"]],
      0.55,1.4,8.2,[3.2,1.7,1.5,1.8],fs=10.5,rowh=0.5,hl_rows=[4])
rect(s,9.1,1.4,3.65,2.9,NAVY)
txt(s,9.32,1.62,3.25,2.6,[[("O ponto mais forte da mesa",14,MINT,True)],[("",6,WHITE,False)],[("A GNLink entrega R$0,78/m³ nas plantas, idêntico ao R$0,80/m³ que a Edge usa no próprio modelo.",11.5,WHITE,False)],[("",5,WHITE,False)],[("A Edge não pode contestar nossa margem sem contestar a dela.",11.5,MINT,True)]])
txt(s,0.55,4.65,12,1.0,[[("A curva cheia entrega R$0,92/m³, acima da premissa, porque os projetos incrementais têm margem superior à base (Eneva 1,29; Argentina 0,91), ambos a capex zero. Capex zero e margem acima da base: o oposto do trade-off normal.",11,DK,False)]])
footer(s,"Modelo da companhia; premissa de off-grid consolidada do equity research (Itaú BBA)")

pg+=1; divider("03","Valor, sinergias e riscos",pg)

# ---- 13 FOOTBALL ----
pg+=1; s=content_slide("Valor","Três metodologias independentes situam o valor em R$1,0 a 1,9 bilhão",title_size=14,pagenum=pg)
txt(s,0.55,1.25,12,0.3,[[("Referência de valor, R$ bilhões",11,GRAY,False)]])
ax0,ax1=3.3,12.5; MAXV=2.0
def xat(v): return ax0+(v/MAXV)*(ax1-ax0)
rows=[("Custo de construir","reproduzir a plataforma",0.65,0.85,GRAY),
      ("Múltiplo descontado","EBITDA 2028, 10x a 12x",1.10,1.65,SAGE),
      ("Curva cheia","EBITDA maduro, múltiplo",1.17,1.75,TEALD),
      ("Valor destravado na Edge","régua do equity research",0.92,1.10,NAVY)]
zx0,zx1=xat(1.4),xat(1.6); rect(s,zx0,1.78,zx1-zx0,3.3,PEACH)
txt(s,zx0-0.7,1.5,2.0,0.25,[[("alvo R$1,4 a 1,6 bi",9.5,DKOR,True)]],align=PP_ALIGN.CENTER)
for i,(lab,sub,lo,hi,col) in enumerate(rows):
    y=1.95+i*0.78
    txt(s,0.55,y-0.03,2.6,0.5,[[(lab,11,NAVY,True)],[(sub,8.5,GRAY,False)]])
    rect(s,ax0,y+0.02,ax1-ax0,0.34,OFF); bx0,bx1=xat(lo),xat(hi); rect(s,bx0,y+0.02,bx1-bx0,0.34,col)
    txt(s,bx0,y+0.05,bx1-bx0,0.28,[[("%.1f a %.1f"%(lo,hi),9.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for v in [0,0.5,1.0,1.5,2.0]:
    xx=xat(v); rect(s,xx,5.12,0.01,0.08,GRAY); txt(s,xx-0.4,5.2,0.8,0.25,[[("R$%.1f"%v,9,GRAY,False)]],align=PP_ALIGN.CENTER)
for i,(h,b) in enumerate([("Piso: custo de construir","Cerca de R$772 mi para reproduzir 686 mil m³/d à eficiência da Edge, e levaria anos."),
                          ("Múltiplo do comprador","EBITDA 2028 (R$172 mi) a 10x, descontado; o mercado precifica a Edge a cerca de 12x."),
                          ("Valor destravado","A GNLink desloca cerca de R$1 bi no valuation da Edge, perto de +3,5% no preço-alvo.")]):
    x=0.55+i*4.08; rect(s,x,5.55,3.9,0.95,OFF); rect(s,x,5.55,3.9,0.07,ORANGE if i==2 else SAGE)
    txt(s,x+0.13,5.66,3.7,0.35,[[(h,11,NAVY,True)]]); txt(s,x+0.13,6.0,3.7,0.5,[[(b,9.5,DK,False)]])
footer(s,"Análise Lorinvest; equity research (BofA, Citi, Itaú BBA, BTG, Bradesco BBI) 2026")

# ---- 14 DESLOCAMENTO ----
pg+=1; s=content_slide("Valor","Pela régua do próprio equity research, a GNLink desloca cerca de R$1 bilhão no valuation da Edge",title_size=12.5,pagenum=pg)
txt(s,0.55,1.25,6.3,0.3,[[("Deslocamento na captura de SS-LNG (BTG)",12,NAVY,True)]])
table(s,[["Passo","Valor"],
         ["Captura da Edge sozinha (65%)","2,21 mi m³/d = R$3,12 bi"],
         ["+ GNLink (curva cheia)","+0,686 mi m³/d"],
         ["Captura combinada (85%)","2,90 mi m³/d = R$4,09 bi"],
         ["= Deslocamento estimado","+R$967 mi"]],
      0.55,1.65,6.3,[3.4,3.0],fs=10.5,rowh=0.46,hl_rows=[4])
txt(s,0.55,3.9,6.3,0.9,[[("O BTG diz que expandir a SS-LNG de 0,4 para 3,8 mi m³/d vale R$4,8 bi e modela 65% de captura. A GNLink adiciona 0,686 mi m³/d a essa conta.",10.5,DK,False)]])
txt(s,7.1,1.25,5.6,0.3,[[("Em moeda do acionista da Compass",12,NAVY,True)]])
for i,(n_,l_) in enumerate([("+R$1,34","por ação da Compass"),("+3,5%","sobre o preço-alvo do BTG"),("65% para 85%","captura de SS-LNG combinada")]):
    y=1.7+i*1.15; rect(s,7.1,y,5.6,1.0,OFF); rect(s,7.1,y,0.1,1.0,NAVY)
    txt(s,7.35,y+0.12,5.3,0.55,[[(n_,26,NAVY,True)]]); txt(s,7.35,y+0.62,5.3,0.35,[[(l_,10.5,GRAY,False)]])
txt(s,0.55,5.1,12,1.1,[[("Camadas aditivas, fora dessa conta: ",11,NAVY,True),("aceleração de rampa via regás in-house (R$100 a 200 mi) e refinanciamento da dívida (R$80 a 140 mi). Os R$4,8 bi e os 65% são do BTG; a atribuição de cerca de R$1 bi é da análise, sobre a régua deles.",11,DK,False)]])
footer(s,"BTG Pactual, Initiation of Coverage jun/2026; análise Lorinvest")

# ---- 15 SINERGIAS ----
pg+=1; s=content_slide("Sinergias","A combinação consolida 26 sinergias em 8 categorias, além do EBITDA entregue",title_size=13.5,pagenum=pg)
txt(s,0.55,1.25,7.6,0.3,[[("Sinergias quantificadas de maior valor",12,NAVY,True)]])
table(s,[["Sinergia","Natureza","Valor"],
         ["G&A Holding","overhead fixo absorvido no dia um","R$210 a 360 mi"],
         ["Fábrica de regás","capex de importação evitado","cerca de R$160 mi"],
         ["Aceleração de rampa","implantação mais rápida no cliente","R$100 a 200 mi"],
         ["Refinanciamento","dívida ao custo de capital da Compass","R$80 a 140 mi"]],
      0.55,1.65,7.6,[2.0,3.4,1.6],fs=10,rowh=0.46)
rect(s,0.55,4.05,7.6,2.1,HLBOX); rect(s,0.55,4.05,0.09,2.1,ORANGE)
txt(s,0.78,4.2,7.2,0.4,[[("Redes locais: valor que vai além da Edge",12.5,NAVY,True)]])
txt(s,0.78,4.62,7.25,1.5,[[("A GNLink leva o gás à cidade e a distribuidora faz a malha interna. Isso expande a rede das 7 distribuidoras da Compass. Todo o potencial de rede local do Sul (433 mil m³/d) está em concessões do grupo; com o Nordeste, 916 mil m³/d.",11,DK,False)]])
txt(s,8.4,1.25,4.3,0.3,[[("As 8 categorias",12,NAVY,True)]])
cats=["Comercial","Mercado e geografia","Operacional","Molécula e arbitragem","Financeiro e capital","CapEx e regás","Cross-segmento","Novos projetos"]
for i,c in enumerate(cats):
    y=1.7+i*0.55; rect(s,8.4,y,4.35,0.45,OFF); rect(s,8.4,y,0.07,0.45,SAGE)
    txt(s,8.6,y,4.1,0.45,[[(c,11,NAVY,True)]],anchor=MSO_ANCHOR.MIDDLE)
footer(s,"Análise Lorinvest; ata RI Compass; equity research 2026")

# ---- 16 RISCOS ----
pg+=1; s=content_slide("Riscos","Seis conflitos exigem tratamento em due diligence e governança antes de abrir preço",title_size=13.5,pagenum=pg)
table(s,[["#","Ponto","Natureza"],
 ["X1","Combinar os dois maiores de GNL off-grid pode levantar questão concorrencial","Antitruste (CADE)"],
 ["X2","Passivos herdados: arbitragem Tradener, dependência de energia e de gás (out/26)","Due diligence"],
 ["X3","Minoritários nas Commit (Copel, governo, Mitsui) para usar o canal de distribuidoras","Governança"],
 ["X4","Cláusulas de change-of-control em contratos (take Eneva, offtakes de 10 a 20 anos)","Contratual"],
 ["X5","O tempo não é neutro: dívida curta a refinanciar e o PR encerra em mar/2036","Timing"],
 ["X6","Arbitragem Tradener de preço e ToP em aberto, distinta do fim de contrato","Due diligence"]],
 0.55,1.3,12.2,[0.7,8.3,3.2],fs=10.5,rowh=0.62)
footer(s,"Histórico GNLink; ata RI Compass jan/2026; análise Lorinvest")

# ---- 17 TATICA ----
pg+=1; s=content_slide("Recomendação","Encaminhar a arbitragem e contratar a Argentina antes de abrir preço, com alvo de R$1,4 a 1,6 bilhão",title_size=12,pagenum=pg)
txt(s,0.55,1.25,7.4,0.35,[[("Sequenciamento antes de abrir preço",13,NAVY,True)]])
for i,(n_,b) in enumerate([("1","Encaminhar a arbitragem Tradener de preço e ToP, que contamina o bloco de vida finita."),
                           ("2","Contratar a Argentina e apresentá-la como ativo central, não como estudo."),
                           ("3","Alinhar a Copa Energia (sócia de GLP), pré-condição antes da abertura."),
                           ("4","Fechar a data da Argentina Fase 2, que antecipa a curva e sobe a âncora.")]):
    y=1.72+i*0.86; rect(s,0.55,y,0.55,0.6,NAVY); txt(s,0.55,y,0.55,0.6,[[(n_,20,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.28,y,6.7,0.62,[[(b,11,DK,False)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,8.3,1.25,4.45,3.05,NAVY)
txt(s,8.55,1.42,4.0,0.4,[[("Zona de negociação",13,MINT,True)]])
for i,(k,v) in enumerate([("Piso (custo de construir)","~R$772 mi"),("Referência (EBITDA 2028 × 10x)","R$1,37 bi"),("Teto (12x)","R$1,65 a 1,75 bi"),("Alvo","R$1,4 a 1,6 bi")]):
    y=1.95+i*0.56
    txt(s,8.55,y,2.7,0.4,[[(k,10.5,WHITE,False)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,11.05,y,1.6,0.4,[[(v,11.5,MINT if i==3 else WHITE,True)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.55,4.85,12.2,1.3,[[("O que sustenta o alvo no terço superior é a alternativa de não vender: as sinergias de molécula (despacho ótimo, arbitragem, backup e térmico) são capturáveis por contrato de swap, sem troca de controle. Essa alternativa precisa estar viva na mesa.",11,DK,False)]])
footer(s,"Análise Lorinvest")

# ---- 18 APENDICE ----
pg+=1; s=content_slide("Apêndice","Glossário, perfis espelhados e fontes",title_size=14,pagenum=pg)
txt(s,0.55,1.25,3.7,0.3,[[("Glossário",12,NAVY,True)]])
gl=[("GNL","gás natural liquefeito, por caminhão ou barcaça"),("GNC","gás natural comprimido, sem regás no cliente"),
    ("Off-grid","fora da rede de gasoduto"),("Regás","regaseificação, converte GNL de volta a gás"),
    ("ToP","take-or-pay, volume mínimo contratado"),("SS-LNG","small scale LNG, o off-grid da Edge"),("CDL","rede local de distribuição")]
y=1.65
for term,defn in gl:
    txt(s,0.55,y,3.75,0.5,[[(term+": ",9.5,NAVY,True),(defn,9.5,DK,False)]]); y+=0.62
txt(s,4.6,1.25,4.0,0.3,[[("Perfis espelhados",12,NAVY,True)]])
table(s,[["","Edge","GNLink"],
         ["Modelo","GNL importado","liquefação nacional"],
         ["Ativo","terminal Santos","3 plantas e regás"],
         ["Geografia","Sudeste","NE, Sul, Argentina"],
         ["Off-grid hoje","150 mil m³/d","280 mil m³/d"],
         ["Capital","balanço Compass","alavancada"]],
      4.6,1.65,4.1,[1.4,1.35,1.35],fs=9,rowh=0.42)
txt(s,9.0,1.25,3.7,0.3,[[("Fontes",12,NAVY,True)]])
srcs=["Histórico GNLink 2026","Investment Presentation 2025","Mapeamento Comercial 2025","Ata RI Compass jan/2026","BofA IoC jun/2026","Citi IoC jun/2026","Itaú BBA IoC jun/2026","BTG Pactual IoC jun/2026","Bradesco BBI IoC jun/2026","Roadshow IPO Compass abr/2026","Análise Lorinvest"]
y=1.65
for sname in srcs:
    txt(s,9.0,y,3.7,0.35,[[("· "+sname,9.5,DK,False)]]); y+=0.42
footer(s,"Compilação Lorinvest")

prs.slides.add_slide(CLO)
prs.save(OUT)
print("OK, slides:",len(prs.slides))
