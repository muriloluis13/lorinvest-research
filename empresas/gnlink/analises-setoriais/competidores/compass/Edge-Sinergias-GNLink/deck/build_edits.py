# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
NAVY=RGBColor(0x00,0x2D,0x5C); MED=RGBColor(0x28,0x47,0x7E); ORANGE=RGBColor(0xEF,0x76,0x22); DKOR=RGBColor(0xC5,0x5A,0x17)
SAGE=RGBColor(0x77,0xA2,0x8A); GRAY=RGBColor(0x7F,0x7F,0x7F); WHITE=RGBColor(0xFF,0xFF,0xFF); TEALD=RGBColor(0x5A,0x9B,0x87)
OFF=RGBColor(0xF2,0xF4,0xF7); DK=RGBColor(0x1E,0x1E,0x1E); LINE=RGBColor(0xD9,0xE0,0xE8); HLROW=RGBColor(0xEA,0xF1,0xEC); RED=RGBColor(0xB9,0x43,0x3A); PEACH=RGBColor(0xFB,0xE7,0xD6)
prs=Presentation("GNLink-Edge.pptx")
def kill_bul(p):
    pPr=p._p.get_or_add_pPr(); pPr.set("marL","0"); pPr.set("indent","0")
    for e in pPr.findall(qn("a:buNone"))+pPr.findall(qn("a:buChar"))+pPr.findall(qn("a:buAutoNum")): pPr.remove(e)
    pPr.append(etree.SubElement(pPr,qn("a:buNone")))
def txt(slide,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=tf.margin_right=Inches(0.02); tf.margin_top=tf.margin_bottom=0; tf.vertical_anchor=anchor
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align; kill_bul(p)
        for (s,sz,c,b) in line:
            r=p.add_run(); r.text=s; r.font.name="Calibri"; r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=b
    return tb
def rect(slide,l,t,w,h,color):
    sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; sp.text_frame.paragraphs[0].text=""; return sp
NOSTYLE="{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
def set_border(cell):
    tcPr=cell._tc.get_or_add_tcPr(); tag=qn("a:lnB")
    for ex in tcPr.findall(tag): tcPr.remove(ex)
    ln=etree.SubElement(tcPr,tag); ln.set("w","6350"); fill=etree.SubElement(ln,qn("a:solidFill")); clr=etree.SubElement(fill,qn("a:srgbClr")); clr.set("val","D9E0E8")
def table(slide,data,l,t,w,colw,fs=10,hl_rows=None,rowh=0.44):
    rows=len(data); cols=len(data[0]); g=slide.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(w),Inches(rowh*rows)); tbl=g.table
    tbl.first_row=False; tbl.horz_banding=False; tbl._tbl.tblPr.set("firstRow","0"); tbl._tbl.tblPr.set("bandRow","0")
    st=tbl._tbl.tblPr.find(qn("a:tableStyleId"))
    if st is None: st=etree.SubElement(tbl._tbl.tblPr,qn("a:tableStyleId"))
    st.text=NOSTYLE
    tot=sum(colw)
    for j,cw in enumerate(colw): tbl.columns[j].width=Emu(int(Inches(w)*cw/tot))
    hl_rows=hl_rows or []
    for i,row in enumerate(data):
        tbl.rows[i].height=Inches(rowh)
        for j,val in enumerate(row):
            cell=tbl.cell(i,j); cell.margin_left=Inches(0.07); cell.margin_right=Inches(0.05); cell.margin_top=cell.margin_bottom=Inches(0.02); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            ishead=i==0
            cell.fill.solid(); cell.fill.fore_color.rgb=NAVY if ishead else (HLROW if i in hl_rows else (WHITE if i%2==1 else OFF))
            set_border(cell); tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; kill_bul(p); p.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
            bold=ishead or (i in hl_rows); r=p.add_run(); r.text=str(val); r.font.name="Calibri"; r.font.size=Pt(fs-0.5 if ishead else fs); r.font.bold=bold; r.font.color.rgb=WHITE if ishead else (NAVY if bold else DK)
    return g
def find_table(s):
    for sh in s.shapes:
        if sh.has_table: return sh
def del_shape(sh): sh._element.getparent().remove(sh._element)
def get14(s):
    for sh in s.shapes:
        try:
            if sh.placeholder_format.idx==14: return sh
        except: pass
def clear_body(s):
    for sh in list(s.shapes):
        keep=False
        try:
            if sh.placeholder_format.idx in (13,14): keep=True
        except: pass
        if sh.has_text_frame and sh.top and sh.top/914400>7.0: keep=True
        if not keep: del_shape(sh)
def repl_all(a,b):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs: r.text=r.text.replace(a,b)
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        for p in c.text_frame.paragraphs:
                            for r in p.runs: r.text=r.text.replace(a,b)
def repl_slide(idx,pairs):
    s=prs.slides[idx]
    for sh in s.shapes:
        frames=[]
        if sh.has_text_frame: frames.append(sh.text_frame)
        if sh.has_table:
            for row in sh.table.rows:
                for c in row.cells: frames.append(c.text_frame)
        for tf in frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    for a,b in pairs: r.text=r.text.replace(a,b)

# ===== GLOBAL =====
repl_all("686,4","690"); repl_all("686","690"); repl_all("Análise Setorial","Tese de combinação")

# ===== Slide 2 (exec) =====
repl_slide(1,[("R$200 mi","R$204 mi"),("R$1,4 bi","R$1,2 bi"),("R$1,4 a 1,6 bilhão","R$1,2 a 1,4 bilhão")])

# ===== Slide 5: CapEx row =====
s=prs.slides[4]; del_shape(find_table(s))
table(s,[["Ano","2027","2028","2029","2030"],
         ["Volume (mi m³/d)","0,15","1,37","2,58","3,80"],
         ["CapEx Fase 2 (R$ mi)","600","450","450","—"],
         ["Margem (R$/m³)","0,80","","","0,86"],
         ["Lucro bruto (R$ mi)","44","410","792","1.195"]],
      0.55,1.7,6.1,[2.1,1,1,1,1],fs=10,rowh=0.38,hl_rows=[2])

# ===== Slide 8: remove G&A =====
s=prs.slides[7]; del_shape(find_table(s))
table(s,[["Bloco","Volume (mil m³/d)","EBITDA (R$ mi)","CapEx incremental"],
         ["3 plantas em operação","280,4","53,7","R$349,9 mi (afundado)"],
         ["+ Projeto Eneva (2027)","+106","+50","zero"],
         ["+ Argentina Fase 1 (1T28)","+150","+50","zero"],
         ["+ Argentina Fase 2 (2030)","+150","+50","zero"],
         ["Total","690","203,7","zero incremental"]],
      0.55,1.3,8.0,[3.0,1.6,1.4,2.0],fs=10.5,rowh=0.5,hl_rows=[5])
for sh in s.shapes:
    if sh.has_text_frame:
        t=sh.text_frame.text
        if t.startswith("Sinergia na veia"):
            sh.text_frame.paragraphs[0].runs[0].text="Base já entregue"
        elif t.startswith("O G&A de R$26"):
            sh.text_frame.paragraphs[0].runs[0].text="As 3 plantas já operam, com R$349,9 mi de capex afundado e 79% da capacidade sob contrato."
repl_slide(7,[("R$230","R$204")])

# ===== Slide 10: chart net =====
s=prs.slides[9]; clear_body(s)
txt(s,0.55,1.25,12,0.3,[[("Trajetória de EBITDA, R$ mi",11,GRAY,False)]])
years=[("2026",-35.8,"-36"),("2027",71.7,"72"),("2028",146.3,"146"),("2029",153.7,"154"),("Maduro",203.7,"204")]
base_y=4.85; maxv=203.7; unit=2.7/maxv; cx=1.2
for i,(yr,v,disp) in enumerate(years):
    x=cx+i*2.25; h=abs(v)*unit; col=RED if v<0 else (NAVY if yr=="Maduro" else SAGE)
    if v>=0:
        rect(s,x,base_y-h,1.5,h,col); txt(s,x-0.1,base_y-h-0.3,1.7,0.3,[[("R$"+disp,12,NAVY,True)]],align=PP_ALIGN.CENTER)
    else:
        rect(s,x,base_y,1.5,h,col); txt(s,x-0.1,base_y+h+0.02,1.7,0.3,[[("R$"+disp,11,RED,True)]],align=PP_ALIGN.CENTER)
    txt(s,x-0.1,base_y+0.42,1.7,0.3,[[(yr,11,GRAY,False)]],align=PP_ALIGN.CENTER)
rect(s,1.0,base_y,11.3,0.02,GRAY)
for i,(h,b) in enumerate([("Vira positivo em 2027","As três plantas atingem EBITDA positivo por conta própria; a Eneva acelera a virada."),
                          ("79% já contratado","222 de 280 mil m³/d sob contrato; a maior parte do salto não depende de novos clientes."),
                          ("Crescimento a capex zero","Eneva e Argentina somam R$150 mi de EBITDA na maturidade, sem capital empregado.")]):
    x=0.55+i*4.08; rect(s,x,5.6,3.9,0.95,OFF); rect(s,x,5.6,3.9,0.07,ORANGE if i==2 else SAGE)
    txt(s,x+0.13,5.71,3.7,0.35,[[(h,11,NAVY,True)]]); txt(s,x+0.13,6.05,3.7,0.5,[[(b,9.5,DK,False)]])
repl_slide(9,[("R$230","R$204")])

# ===== Slide 11: margin-only =====
s=prs.slides[10]; del_shape(find_table(s))
table(s,[["","Volume","Margem (R$/m³)"],
         ["3 plantas plenas (2029)","280,4 mil m³/d","0,78"],
         ["+ Eneva","106 mil m³/d","1,29"],
         ["+ Argentina F1 e F2","300 mil m³/d","0,91 cada"],
         ["Curva cheia (média ponderada)","690 mil m³/d","0,92"],
         ["Premissa off-grid da própria Edge","","0,80 a 0,86"]],
      0.55,1.4,8.2,[3.6,2.4,2.2],fs=10.5,rowh=0.5,hl_rows=[4])

# ===== Slide 13: football field recomputed (net) =====
def br(x): return ("%.2f"%x).replace(".",",")
s=prs.slides[12]; clear_body(s)
txt(s,0.55,1.25,12,0.3,[[("Referência de valor, R$ bilhões",11,GRAY,False)]])
ax0,ax1=3.3,12.5; MAXV=2.0
def xat(v): return ax0+(v/MAXV)*(ax1-ax0)
rows=[("Custo de construir","reproduzir a plataforma",0.65,0.85,GRAY),
      ("Múltiplo descontado","EBITDA 2028, 10x a 12x",1.17,1.40,SAGE),
      ("Curva cheia","EBITDA maduro, múltiplo",1.30,1.56,TEALD),
      ("Valor destravado na Edge","régua do equity research",0.92,1.10,NAVY)]
zx0,zx1=xat(1.2),xat(1.4); rect(s,zx0,1.78,zx1-zx0,3.3,PEACH)
txt(s,zx0-0.7,1.5,2.0,0.25,[[("alvo R$1,2 a 1,4 bi",9.5,DKOR,True)]],align=PP_ALIGN.CENTER)
for i,(lab,sub,lo,hi,col) in enumerate(rows):
    y=1.95+i*0.78
    txt(s,0.55,y-0.03,2.6,0.5,[[(lab,11,NAVY,True)],[(sub,8.5,GRAY,False)]])
    rect(s,ax0,y+0.02,ax1-ax0,0.34,OFF); bx0,bx1=xat(lo),xat(hi); rect(s,bx0,y+0.02,bx1-bx0,0.34,col)
    txt(s,bx0,y+0.05,bx1-bx0,0.28,[[(br(lo)+" a "+br(hi),9.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
for v in [0,0.5,1.0,1.5,2.0]:
    xx=xat(v); rect(s,xx,5.12,0.01,0.08,GRAY); txt(s,xx-0.4,5.2,0.8,0.25,[[("R$"+("%.1f"%v).replace(".",","),9,GRAY,False)]],align=PP_ALIGN.CENTER)
for i,(h,b) in enumerate([("Piso: custo de construir","Cerca de R$772 mi para reproduzir 690 mil m³/d à eficiência da Edge, e levaria anos."),
                          ("Múltiplo do comprador","EBITDA 2028 (R$146 mi) a 10x, descontado; o mercado precifica a Edge a cerca de 12x."),
                          ("Valor destravado","A GNLink desloca cerca de R$1 bi no valuation da Edge, perto de +3,5% no preço-alvo.")]):
    x=0.55+i*4.08; rect(s,x,5.55,3.9,0.95,OFF); rect(s,x,5.55,3.9,0.07,ORANGE if i==2 else SAGE)
    txt(s,x+0.13,5.66,3.7,0.35,[[(h,11,NAVY,True)]]); txt(s,x+0.13,6.0,3.7,0.5,[[(b,9.5,DK,False)]])
repl_slide(12,[("R$1,0 a 1,9 bilhão","R$0,8 a 1,6 bilhão")])

# ===== Slide 14: deslocamento with 0,690 =====
repl_slide(13,[("R$967","R$973"),("+R$1,34","+R$1,35")])

# ===== Slide 15: remove G&A row =====
s=prs.slides[14]; del_shape(find_table(s))
table(s,[["Sinergia","Natureza","Valor"],
         ["Fábrica de regás","capex de importação evitado","~R$160 mi"],
         ["Aceleração de rampa","implantação mais rápida no cliente","R$100 a 200 mi"],
         ["Refinanciamento","dívida ao custo de capital da Compass","R$80 a 140 mi"]],
      0.55,1.65,7.6,[2.0,3.4,1.6],fs=10,rowh=0.46)
repl_slide(14,[("26 sinergias","25 sinergias")])

# ===== Slide 17: zona =====
repl_slide(16,[("R$1,65 a 1,75 bi","R$1,4 bi"),("R$1,37 bi","R$1,2 bi"),("R$1,4 a 1,6 bi","R$1,2 a 1,4 bi")])

prs.save("GNLink-Edge.pptx")
print("OK salvo, slides:",len(prs.slides))
